#!/usr/bin/env python3
"""
Investigation Agent Demo Implementation

This module provides a demonstration of the Investigation Agent system's capabilities
for analyzing metric anomalies in social media data. It implements a simplified version
of the full system with mock data for demonstration purposes.

Key Features:
    - Natural language question processing
    - Multi-metric anomaly detection
    - Hypothesis generation and validation
    - Automated visualization creation
    - Report generation (PowerPoint and text)

Example Usage:
    ```python
    # Basic usage
    agent = DemoInvestigationAgent()
    results = agent.investigate("Investigate the spike in engagement during April 2025")
    
    # Access specific components
    metrics = agent._analyze_metrics()
    hypotheses = agent._generate_hypotheses(metrics)
    ```

Classes:
    - DemoInvestigationAgent: Main class implementing the investigation logic
"""

import sys
import time
from datetime import datetime, timedelta
import pandas as pd
import numpy as np
from pathlib import Path
import plotly.express as px
from pptx import Presentation
import spacy
from tqdm import tqdm
from pptx.util import Inches, Pt
import matplotlib.pyplot as plt
import json
from typing import Dict, List, Tuple, Any

class DemoInvestigationAgent:
    """
    A demonstration implementation of the Investigation Agent system.
    
    This class provides a simplified version of the full system's capabilities,
    using mock data and basic implementations of core features.
    
    Attributes:
        nlp (spacy.language): SpaCy language model for text processing
        metrics (dict): Mock metric data for demonstration
        weights (dict): Investigation component weights
    """
    
    def __init__(self):
        """
        Initialize the demo agent with mock data and models.
        
        Sets up:
            - SpaCy NLP model
            - Mock metric data
            - Investigation weights
            - Output directory
        """
        # Initialize NLP
        self.nlp = spacy.load("en_core_web_sm")
        
        # Setup mock data
        self._setup_mock_data()
        
        # Initialize weights
        self.weights = {
            'selection': 0.50,
            'execution': 0.20,
            'context': 0.20,
            'temporal': 0.05,
            'relationship': 0.05
        }
        
        # Ensure output directory exists
        Path("demo_output").mkdir(exist_ok=True)

    def _setup_mock_data(self):
        """
        Initialize mock metric data for demonstration purposes.
        
        Creates synthetic time series data for:
            - Engagement metrics
            - Comment counts
            - Share counts
        """
        # Generate date range
        dates = pd.date_range(start='2025-01-01', end='2025-12-31', freq='D')
        
        # Base metrics with seasonal pattern and trend
        base = np.sin(np.arange(len(dates)) * 2 * np.pi / 365) * 100 + \
               np.random.normal(0, 10, len(dates))
        
        # Add trend and anomalies
        trend = np.linspace(0, 50, len(dates))
        anomalies = np.zeros(len(dates))
        anomalies[100:130] = 100  # Spike in April-May
        
        # Create individual metrics with variations
        self.metrics = {
            'engagement': base + trend + anomalies + np.random.normal(0, 20, len(dates)),
            'comments': base * 0.5 + trend * 0.3 + anomalies * 0.7 + np.random.normal(0, 15, len(dates)),
            'shares': base * 0.3 + trend * 0.2 + anomalies * 0.5 + np.random.normal(0, 10, len(dates))
        }
        
        # Add dates
        self.dates = dates

    def investigate(self, question: str) -> Dict[str, Any]:
        """
        Main investigation method that processes a question and generates analysis.
        
        Args:
            question (str): The investigation question (e.g., "Investigate the spike 
                          in engagement during April 2025")
        
        Returns:
            dict: Investigation results including metrics, hypotheses, and visualizations
        
        Example:
            ```python
            agent = DemoInvestigationAgent()
            results = agent.investigate("Why did engagement spike in April 2025?")
            ```
        """
        # Process the question
        print("\n🔍 Starting Investigation...")
        print(f"Question: {question}")
        
        # Extract entities and time range
        print("\n1️⃣ Analyzing question...")
        doc = self.nlp(question)
        entities = [(ent.text, ent.label_) for ent in doc.ents]
        print(f"Found entities: {entities}")
        
        # Analyze metrics
        print("\n2️⃣ Analyzing metrics...")
        metrics_info = self._analyze_metrics()
        print(f"Metrics analyzed: {', '.join(metrics_info['metrics'])}")
        print(f"Found {len(metrics_info['anomalies'])} anomalies")
        
        # Generate hypotheses
        print("\n3️⃣ Generating hypotheses...")
        hypotheses = self._generate_hypotheses(metrics_info)
        
        # Create visualizations
        self._create_visualizations(metrics_info)
        
        # Generate report
        self._generate_report(question, metrics_info, hypotheses)
        
        return {
            'entities': entities,
            'metrics_info': metrics_info,
            'hypotheses': hypotheses
        }

    def _analyze_metrics(self) -> Dict[str, Any]:
        """
        Analyze mock metric data to detect anomalies and patterns.
        
        Returns:
            dict: Analysis results including metrics, anomalies, and trends
        """
        # Calculate basic statistics
        stats = {}
        anomalies = []
        
        for metric, values in self.metrics.items():
            mean = np.mean(values)
            std = np.std(values)
            threshold = mean + 2 * std
            
            # Detect anomalies (simple threshold-based detection)
            anomaly_dates = self.dates[values > threshold]
            if len(anomaly_dates) > 0:
                anomalies.append({
                    'metric': metric,
                    'dates': anomaly_dates.strftime('%Y-%m-%d').tolist(),
                    'severity': 'high'
                })
            
            stats[metric] = {
                'mean': mean,
                'std': std,
                'trend': 'increasing' if values[-1] > values[0] else 'decreasing'
            }
        
        return {
            'metrics': list(self.metrics.keys()),
            'stats': stats,
            'anomalies': anomalies,
            'trends': {
                'overall': 'positive',
                'seasonal': True
            }
        }

    def _generate_hypotheses(self, metrics_info: Dict[str, Any]) -> List[Dict[str, Any]]:
        """
        Generate hypotheses based on metric analysis.
        
        Args:
            metrics_info (dict): Metrics analysis results
        
        Returns:
            list: Generated hypotheses with confidence scores
        """
        # Demo hypotheses
        hypotheses = [
            {
                'description': 'Significant increase in political content engagement during election period',
                'confidence': 0.85,
                'evidence': ['temporal_pattern', 'metric_correlation'],
                'impact': 'high'
            },
            {
                'description': 'Viral content cascade effect leading to temporary engagement spike',
                'confidence': 0.75,
                'evidence': ['rapid_growth', 'cross_platform'],
                'impact': 'medium'
            },
            {
                'description': 'Seasonal factors combined with trending topics',
                'confidence': 0.65,
                'evidence': ['seasonal_pattern', 'topic_analysis'],
                'impact': 'medium'
            }
        ]
        
        return hypotheses

    def _create_visualizations(self, metrics_info: Dict[str, Any]):
        """
        Create and save visualizations of the analysis results.
        
        Args:
            metrics_info (dict): Metrics analysis results
        """
        # Create metric trends plot
        plt.figure(figsize=(12, 6))
        for metric, values in self.metrics.items():
            plt.plot(self.dates, values, label=metric)
        plt.title('Metric Trends Over Time')
        plt.xlabel('Date')
        plt.ylabel('Value')
        plt.legend()
        plt.tight_layout()
        plt.savefig('demo_output/metric_trends.png')
        plt.close()
        
        # Create anomalies plot
        plt.figure(figsize=(12, 6))
        for metric, values in self.metrics.items():
            mean = np.mean(values)
            std = np.std(values)
            threshold = mean + 2 * std
            anomalies = values > threshold
            
            plt.plot(self.dates, values, label=f'{metric} (normal)')
            plt.scatter(self.dates[anomalies], values[anomalies], 
                       label=f'{metric} (anomaly)', marker='x')
        
        plt.title('Detected Anomalies')
        plt.xlabel('Date')
        plt.ylabel('Value')
        plt.legend()
        plt.tight_layout()
        plt.savefig('demo_output/anomalies.png')
        plt.close()

    def _generate_report(self, question: str, metrics_info: Dict[str, Any], 
                        hypotheses: List[Dict[str, Any]]):
        """
        Generate a PowerPoint report summarizing the investigation.
        
        Args:
            question (str): Original investigation question
            metrics_info (dict): Metrics analysis results
            hypotheses (list): Generated hypotheses
        """
        prs = Presentation()
        
        # Title slide
        title_slide = prs.slides.add_slide(prs.slide_layouts[0])
        title_slide.shapes.title.text = "Investigation Report"
        title_slide.placeholders[1].text = f"Analysis of: {question}"
        
        # Metrics summary slide
        metrics_slide = prs.slides.add_slide(prs.slide_layouts[1])
        metrics_slide.shapes.title.text = "Metrics Analysis"
        metrics_content = metrics_slide.placeholders[1]
        metrics_text = f"Analyzed Metrics: {', '.join(metrics_info['metrics'])}\n"
        metrics_text += f"Anomalies Found: {len(metrics_info['anomalies'])}\n"
        metrics_text += f"Overall Trend: {metrics_info['trends']['overall']}"
        metrics_content.text = metrics_text
        
        # Hypotheses slide
        hyp_slide = prs.slides.add_slide(prs.slide_layouts[1])
        hyp_slide.shapes.title.text = "Generated Hypotheses"
        hyp_content = hyp_slide.placeholders[1]
        hyp_text = ""
        for idx, hyp in enumerate(hypotheses, 1):
            hyp_text += f"{idx}. {hyp['description']}\n"
            hyp_text += f"   Confidence: {hyp['confidence']:.0%}\n"
        hyp_content.text = hyp_text
        
        # Visualizations slide
        vis_slide = prs.slides.add_slide(prs.slide_layouts[1])
        vis_slide.shapes.title.text = "Visualizations"
        
        # Add metric trends
        trends_img = vis_slide.shapes.add_picture(
            'demo_output/metric_trends.png',
            Inches(1), Inches(2),
            height=Inches(3)
        )
        
        # Save presentation
        prs.save('demo_output/investigation_report.pptx')

def main():
    if len(sys.argv) < 2:
        print("Usage: python demo.py 'Your investigation question here'")
        sys.exit(1)
    
    question = sys.argv[1]
    agent = DemoInvestigationAgent()
    results = agent.investigate(question)

if __name__ == "__main__":
    main() 